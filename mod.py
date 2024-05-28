from pptx import Presentation

placeholder = """Evaluation only.
Created with Aspose.Slides for .NET Standard 2.0 23.12.
Copyright 2004-2023Aspose Pty Ltd."""

placeholder = """Evaluation only.
Created with Aspose.Slides for Python via .NET 24.5.
Copyright 2004-2024Aspose Pty Ltd."""

def modIt(pptx_file_path):
    presentation = Presentation(pptx_file_path)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text == placeholder:
                # Delete the shape
                sp = shape._element
                sp.getparent().remove(sp)

    presentation.save(pptx_file_path)